import tkinter as tk
from tkinter import filedialog, messagebox
import customtkinter
from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
import os

def selecionar_arquivo():
    caminho_arquivo = filedialog.askopenfilename(filetypes=[("PowerPoint files", "*.pptx")])
    if caminho_arquivo:
        label_status.configure(text=f"Arquivo Selecionado: {caminho_arquivo}")
        converter_para_pdf(caminho_arquivo)

def converter_para_pdf(caminho_arquivo):
    try:
        # Perguntar ao usuário onde salvar o PDF
        caminho_pdf = filedialog.asksaveasfilename(defaultextension=".pdf",
                                                   filetypes=[("PDF files", "*.pdf")],
                                                   title="Salvar PDF como")
        if not caminho_pdf:
            label_status.configure(text="Conversão cancelada pelo usuário.")
            return

        # Carregar apresentação PPTX
        prs = Presentation(caminho_arquivo)

        # Criar o PDF usando reportlab
        c = canvas.Canvas(caminho_pdf, pagesize=letter)

        num_slides = len(prs.slides)
        slide_count = 0

        barra_progresso.pack(pady=10)
        barra_progresso.set(0)
        label_porcentagem.configure(text="0%")

        for slide in prs.slides:
            text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
            c.drawString(100, 750, text)
            c.showPage()

            # Atualizar a barra de progresso
            slide_count += 1
            progress = (slide_count / num_slides) * 100
            barra_progresso.set(progress / 100)
            label_porcentagem.configure(text=f"{int(progress)}%")
            janela.update_idletasks()

        c.save()

        label_status.configure(text=f"Conversão concluída! PDF salvo como {caminho_pdf}")
        messagebox.showinfo("Sucesso", f"Conversão concluída! PDF salvo como {caminho_pdf}")

    except Exception as e:
        label_status.configure(text=f"Erro ao converter arquivo: {str(e)}")
        messagebox.showerror("Erro", f"Erro ao converter arquivo: {str(e)}")

customtkinter.set_appearance_mode("System")
customtkinter.set_default_color_theme("dark-blue")

janela = customtkinter.CTk()
janela.geometry("720x480")
janela.title("Conversor PPTX para PDF")

label_titulo = customtkinter.CTkLabel(janela, text="Selecione um arquivo PPTX para converter", font=("Consolas bold", 17))
label_titulo.pack(padx=10, pady=10)

botao_selecionar = customtkinter.CTkButton(janela, text="Selecionar Arquivo", command=selecionar_arquivo, width=250)
botao_selecionar.pack(pady=10)

label_status = customtkinter.CTkLabel(janela, text="")
label_status.pack(pady=10)

label_porcentagem = customtkinter.CTkLabel(janela, text="")
label_porcentagem.pack(pady=10)

barra_progresso = customtkinter.CTkProgressBar(janela, width=400)
barra_progresso.set(0)
barra_progresso.pack_forget()

janela.mainloop()
